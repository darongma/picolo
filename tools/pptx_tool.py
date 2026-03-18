"""
Read, create, and edit PowerPoint (.pptx) presentations.
"""
import pptx
from pptx.util import Inches, Pt

# Read tool
tool_spec_read = {
    "name": "pptx_read",
    "description": "Extract text from a PowerPoint presentation (.pptx). Optionally specify slide numbers (1-indexed).",
    "parameters": {
        "type": "object",
        "properties": {
            "file_path": {
                "type": "string",
                "description": "Path to the .pptx file."
            },
            "slides": {
                "type": "array",
                "items": {"type": "integer"},
                "description": "Optional list of slide numbers to extract (e.g., [1,3]). If omitted, all slides are extracted."
            }
        },
        "required": ["file_path"]
    }
}

def _extract_shape_text(shape):
    if not shape.has_text_frame:
        return ""
    return shape.text_frame.text.strip()

def run_read(file_path: str, slides: list = None) -> str:
    try:
        prs = pptx.Presentation(file_path)
        out = []
        slide_iter = range(len(prs.slides))
        if slides:
            slide_iter = [s-1 for s in slides if 1 <= s <= len(prs.slides)]
        for idx in slide_iter:
            slide = prs.slides[idx]
            texts = []
            for shape in slide.shapes:
                txt = _extract_shape_text(shape)
                if txt:
                    texts.append(txt)
            slide_text = "\n".join(texts)
            out.append(f"--- Slide {idx+1} ---\n{slide_text}")
        return "\n".join(out) if out else "(presentation has no extractable text)"
    except Exception as e:
        return f"Error: {e}"

# Create tool
tool_spec_create = {
    "name": "pptx_create",
    "description": "Create a new PowerPoint presentation. Provide a list of slides; each slide is a dict with optional 'title' and 'content' (string or list of bullet points).",
    "parameters": {
        "type": "object",
        "properties": {
            "file_path": {
                "type": "string",
                "description": "Path where the .pptx file will be saved."
            },
            "slides": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "title": {"type": "string"},
                        "content": {"type": ["string", "array"]}
                    }
                },
                "description": "List of slide definitions. Each slide can have a title and/or content."
            }
        },
        "required": ["file_path", "slides"]
    }
}

def _add_slide_with_content(prs, layout_idx, title=None, content=None):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    # Handle title
    if title is not None and slide.shapes.title:
        slide.shapes.title.text = title
    # Handle content
    if content is not None:
        if isinstance(content, str):
            content = [content]
        # Find a placeholder for body (index 1) or add a textbox
        if len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()
            for i, item in enumerate(content):
                if i == 0:
                    tf.text = item
                else:
                    p = tf.add_paragraph()
                    p.text = item
        else:
            # Fallback: add a textbox below title
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(4)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            for i, item in enumerate(content):
                if i == 0:
                    tf.text = item
                else:
                    tf.add_paragraph().text = item
    return slide

def run_create(file_path: str, slides: list) -> str:
    try:
        prs = pptx.Presentation()
        # Use layout 1 (Title and Content) if available, else fallback
        layout_idx = 1 if len(prs.slide_layouts) > 1 else 0
        for s in slides:
            title = s.get('title')
            content = s.get('content')
            _add_slide_with_content(prs, layout_idx, title, content)
        prs.save(file_path)
        return f"Created PowerPoint with {len(slides)} slides at {file_path}"
    except Exception as e:
        return f"Error: {e}"

# Append slide tool
tool_spec_append = {
    "name": "pptx_append_slide",
    "description": "Append one or more new slides to an existing PowerPoint presentation.",
    "parameters": {
        "type": "object",
        "properties": {
            "file_path": {
                "type": "string",
                "description": "Path to the .pptx file."
            },
            "slides": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "title": {"type": "string"},
                        "content": {"type": ["string", "array"]}
                    }
                },
                "description": "List of slide definitions to append."
            }
        },
        "required": ["file_path", "slides"]
    }
}

def run_append(file_path: str, slides: list) -> str:
    try:
        prs = pptx.Presentation(file_path)
        layout_idx = 1 if len(prs.slide_layouts) > 1 else 0
        for s in slides:
            title = s.get('title')
            content = s.get('content')
            _add_slide_with_content(prs, layout_idx, title, content)
        prs.save(file_path)
        return f"Appended {len(slides)} slides to {file_path}"
    except Exception as e:
        return f"Error: {e}"

# Edit slide tool
tool_spec_edit = {
    "name": "pptx_edit_slide",
    "description": "Edit an existing slide in a PowerPoint presentation. Replace title and/or content for the specified slide (1-indexed).",
    "parameters": {
        "type": "object",
        "properties": {
            "file_path": {
                "type": "string",
                "description": "Path to the .pptx file."
            },
            "slide_num": {
                "type": "integer",
                "description": "Slide number to edit (1-indexed)."
            },
            "title": {
                "type": "string",
                "description": "New title text. If omitted, title remains unchanged."
            },
            "content": {
                "type": ["string", "array"],
                "description": "New content; string or list of strings. If omitted, content remains unchanged."
            }
        },
        "required": ["file_path", "slide_num"]
    }
}

def run_edit(file_path: str, slide_num: int, title=None, content=None) -> str:
    try:
        prs = pptx.Presentation(file_path)
        if slide_num < 1 or slide_num > len(prs.slides):
            return f"Error: slide_num {slide_num} out of range (1-{len(prs.slides)})"
        slide = prs.slides[slide_num-1]
        changed = False
        if title is not None and slide.shapes.title:
            slide.shapes.title.text = title
            changed = True
        if content is not None:
            # Try to find a body placeholder
            body_shape = None
            if len(slide.placeholders) > 1:
                body_shape = slide.placeholders[1]
            if body_shape:
                tf = body_shape.text_frame
                tf.clear()
                if isinstance(content, str):
                    tf.text = content
                else:
                    for i, item in enumerate(content):
                        if i == 0:
                            tf.text = item
                        else:
                            p = tf.add_paragraph()
                            p.text = item
                changed = True
            else:
                # No body placeholder; could add a textbox but editing ambiguous; skip
                pass
        if not changed:
            return f"No changes made to slide {slide_num}"
        prs.save(file_path)
        return f"Updated slide {slide_num} in {file_path}"
    except Exception as e:
        return f"Error: {e}"

tool_specs = [tool_spec_read, tool_spec_create, tool_spec_append, tool_spec_edit]

tools = {
    "pptx_read": run_read,
    "pptx_create": run_create,
    "pptx_append_slide": run_append,
    "pptx_edit_slide": run_edit
}
